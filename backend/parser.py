import os
import sys
import time
import hashlib
import fitz  # PyMuPDF
from docx import Document
from pptx import Presentation
import tiktoken
import openai
import concurrent.futures
import gspread
from oauth2client.service_account import ServiceAccountCredentials
from datetime import datetime
import re

# Enable UTF-8 output
sys.stdout.reconfigure(encoding='utf-8')

import base64
import json

creds_json = base64.b64decode(os.getenv("GOOGLE_CREDENTIALS_B64")).decode("utf-8")
with open("credentials.json", "w") as f:
    f.write(creds_json)

scope = ["https://spreadsheets.google.com/feeds", "https://www.googleapis.com/auth/drive"]
creds = ServiceAccountCredentials.from_json_keyfile_name("credentials.json", scope)

# Config
SYSTEM_PROMPT = """You are an expert AI tasked with rigorously evaluating university-level course materials.

You will receive three inputs:

Course Overview: learning outcomes (CLOs), weekly schedule, course objectives.

Lesson File: content from one week (e.g., PowerPoint, transcript).

Rubric: 8 criteria scored 1–4 (1=Poor, 4=Excellent).

Your task:

Step 1: Relevance Check

Identify the week’s topic and targeted CLOs.

Assess how well the lesson content aligns with the week’s topic and CLOs.

Note any missing content, off-topic material, or misalignment in depth.

Summarize your findings concisely.

Step 2: Rubric-Based Evaluation

For each of the following criteria, assign a score from 1 to 4 and provide a brief explanation (1–3 sentences), citing examples or omissions from the lesson:

Clarity and Organization

Relevance and Accuracy

Engagement and Interaction (5E’s: Engage, Explore, Explain, Elaborate, Evaluate)

Depth and Breadth

Assessment Samples and Feedback

Use of Technology

Alignment with CLOs

Student Support Resources

Step 3: Suggestions for Improvement

For every criterion scored 1 or 2, provide 1–2 actionable, pedagogically sound improvement suggestions that are practical and efficient.

Output format (use exactly this layout):
1. Relevance Summary - {ENTER COURSE NAME HERE (FROM FILES)} {ENTER $ AFTER FULL HEADING}

A brief paragraph summarizing lesson alignment with the week’s topic and CLOs, noting any major gaps or issues.

2. Rubric Evaluation

For each criterion output the following block:

makefile
Copy
Edit
Criterion: <Criterion Name>
Score: <1–4>
Explanation: <Brief, evidence-based justification>
Separate each criterion block by a blank line.

3. Suggestions for Improvement
For each criterion scoring 1–3, provide 2–3 focused, actionable, and creative suggestions. Go beyond surface-level fixes — aim to pinpoint specific pages, paragraphs, or sections that need revision.

Suggestions should include:

Targeted rewrites or structural tweaks

Specific content to add, cut, or clarify

Creative ways to boost clarity, flow, or engagement

Exact spots (e.g., “Page 5, para 2”) where improvements will have the most impact

Example format:

Clarity and Organization:

Page 3, para 1: Reframe the opening with a sharper topic sentence to anchor the argument.

Section 2: Add brief transitions between subpoints to improve logical flow.

Engagement and Interaction:

Page 7: Replace jargon with a vivid example or anecdote to draw readers in.

Begin your detailed evaluation now."""
openai.api_key = os.getenv("API_KEY")

CACHE_DIR = ".cache"
os.makedirs(CACHE_DIR, exist_ok=True)

titles = ["Rubric", "Lesson", "Course Overview"]

# Google Sheets Setup
def init_sheets():
    scope = ["https://spreadsheets.google.com/feeds", "https://www.googleapis.com/auth/drive"]
    creds = ServiceAccountCredentials.from_json_keyfile_name("credentials.json", scope)
    client = gspread.authorize(creds)
    sheet = client.open("Lecture Evaluation Log").sheet1
    return sheet

def log_to_sheets(sheet, syllabus_file, lesson_file, rubric_grades, parse_time, openai_time):
    row = [
        datetime.now().strftime("%Y-%m-%d %H:%M:%S"),
        os.path.basename(syllabus_file),
        os.path.basename(lesson_file),
        rubric_grades,
        f"{parse_time:.2f}s",
        f"{openai_time:.2f}s"
    ]
    sheet.append_row(row)

# Utility: hash file content for caching
def hash_file(filepath):
    hasher = hashlib.sha256()
    with open(filepath, 'rb') as f:
        hasher.update(f.read())
    return hasher.hexdigest()

# Parse DOCX
def parse_docx(file_path):
    doc = Document(file_path)
    return "\n".join(
        f"\n[Page {i+1}] {para.text}" if para.style.name.startswith('Heading') else para.text
        for i, para in enumerate(doc.paragraphs) if para.text.strip()
    )

# Parse PDF
def parse_pdf(file_path):
    doc = fitz.open(file_path)
    return "\n\n".join(
        f"[Page {i+1}]\n{page.get_text().strip()}" for i, page in enumerate(doc) if page.get_text().strip()
    )

# Parse PPTX
def parse_pptx(file_path):
    prs = Presentation(file_path)
    output = []
    for idx, slide in enumerate(prs.slides, 1):
        texts = [f"[Slide {idx}]"]
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                texts.append(shape.text.strip())
        output.append("\n".join(texts))
    return "\n\n".join(output)

# Smart parser dispatcher
def parse_file(file_path):
    ext = os.path.splitext(file_path)[1].lower()
    if ext == ".pdf":
        return parse_pdf(file_path)
    elif ext == ".docx":
        return parse_docx(file_path)
    elif ext == ".pptx":
        return parse_pptx(file_path)
    else:
        return "Unsupported file type."

# Count tokens
def count_tokens(text):
    enc = tiktoken.get_encoding("cl100k_base")
    return len(enc.encode(text))

# Cache-aware parse wrapper
def cached_parse(file_path):
    file_hash = hash_file(file_path)
    cache_path = os.path.join(CACHE_DIR, f"{file_hash}.txt")
    if os.path.exists(cache_path):
        with open(cache_path, 'r', encoding='utf-8') as f:
            return f.read()
    parsed = parse_file(file_path)
    with open(cache_path, 'w', encoding='utf-8') as f:
        f.write(parsed)
    return parsed

# Call OpenAI
def call_openai_api(prompt):
    messages = [
        {"role": "system", "content": SYSTEM_PROMPT},
        {"role": "user", "content": prompt}
    ]
    response = openai.chat.completions.create(
        model="gpt-4o-mini",
        messages=messages,
        max_tokens=4096,
        temperature=0.45,
        top_p=1.0,
        frequency_penalty=0.1,
        presence_penalty=0.0
    )
    return response.choices[0].message.content

# Extract rubric section using regex
def extract_rubric_section(response_text):
    match = re.search(r"2\.\s*Rubric Evaluation\s+([\s\S]*?)\s+3\.\s*Suggestions for Improvement", response_text, re.IGNORECASE)
    return match.group(1).strip() if match else "N/A"

# Main
def main(file_paths):
    if len(file_paths) != 3:
        print("Usage: python fast_parser.py <rubric.docx> <lesson.pptx> <overview.pdf>")
        sys.exit(1)

    t0 = time.time()

    print("\n--- Parsing files in parallel ---")
    with concurrent.futures.ThreadPoolExecutor() as executor:
        parsed_results = list(executor.map(cached_parse, file_paths))

    for i, result in enumerate(parsed_results):
        print(f"\n✅ Finished parsing: {titles[i]} ({round(time.time() - t0, 2)}s)")

    combined = "\n\n".join(f"### {titles[i]} ###\n{parsed_results[i]}" for i in range(3))
    token_count = count_tokens(combined)
    parse_duration = time.time() - t0

    print(f"\n🧠 Total Tokens: {token_count}")
    print(f"⏱️ Parsing Duration: {round(parse_duration, 2)}s")

    if token_count > 5000:
        print("⚠️ Too many tokens. Skipping OpenAI call.")
        sys.exit(1)

    print("\n--- Calling OpenAI API ---")
    t1 = time.time()
    try:
        response = call_openai_api(combined)
        api_duration = time.time() - t1
        print(f"\n✅ OpenAI Response (⏱️ {round(api_duration, 2)}s):\n")
        print(response)

        rubric_eval = extract_rubric_section(response)
        sheet = init_sheets()
        log_to_sheets(sheet, file_paths[2], file_paths[1], rubric_eval, parse_duration, api_duration)

    except Exception as e:
        print(f"❌ OpenAI API call failed: {e}")

if __name__ == "__main__":
    main(sys.argv[1:])
